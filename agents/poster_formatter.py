# agents/poster_formatter.py
"""
Agent 3 — Poster Formatter
- Takes summarized sections (with bullets and image refs) and produces:
  - Poster dataclass (structured JSON)
  - Markdown preview (images referenced as markdown image links)
  - Optional simple PPTX export (uses python-pptx if installed)
"""
import os
import json
import logging
from typing import List, Optional
from textwrap import dedent

from . import Poster, Section, ImageRef

logger = logging.getLogger(__name__)

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    PPTX_AVAILABLE = True
except Exception:
    PPTX_AVAILABLE = False

class PosterFormatter:
    def __init__(self, output_dir: str = "data/outputs"):
        self.output_dir = output_dir
        os.makedirs(self.output_dir, exist_ok=True)

    def format_to_poster(self, title: Optional[str], authors: Optional[List[str]], summarized_sections: List[Section], layout: Optional[dict] = None) -> Poster:
        poster = Poster(title=title, authors=authors, sections=summarized_sections, layout=layout or {})
        return poster

    def poster_to_json(self, poster: Poster, out_path: Optional[str] = None) -> str:
        out_path = out_path or os.path.join(self.output_dir, "poster.json")
        def section_to_dict(s: Section):
            return {
                "title": s.title,
                "text": s.text,
                "start_page": s.start_page,
                "end_page": s.end_page,
                "images": [vars(img) for img in s.images],
                "metadata": s.metadata
            }
        payload = {
            "title": poster.title,
            "authors": poster.authors,
            "layout": poster.layout,
            "sections": [section_to_dict(s) for s in poster.sections],
            "metadata": poster.metadata
        }
        with open(out_path, "w", encoding="utf-8") as fh:
            json.dump(payload, fh, indent=2, ensure_ascii=False)
        logger.info("Poster JSON written to %s", out_path)
        return out_path

    def poster_to_markdown(self, poster: Poster) -> str:
        md = []
        md.append(f"# {poster.title or 'Untitled Poster'}\n")
        if poster.authors:
            md.append(", ".join(poster.authors) + "\n")
        for sec in poster.sections:
            # If the section has bullets in metadata (from summarizer), use them. Otherwise, use text.
            md.append(f"## {sec.title}\n")
            bullets = sec.metadata.get("bullets") if sec.metadata.get("bullets") else None
            if bullets:
                for b in bullets:
                    md.append(f"- {b}")
            else:
                # fallback short excerpt
                excerpt = (sec.text[:500] + "...") if sec.text and len(sec.text) > 500 else (sec.text or "")
                if excerpt:
                    md.append(excerpt)
            # images
            if sec.images:
                for img in sec.images:
                    # prefer path if available
                    img_path = img.path or img.id
                    caption = img.caption or ""
                    md.append(f"![{caption}]({img_path})")
            md.append("\n")
        markdown = "\n".join(md)
        return markdown

    def write_markdown_file(self, poster: Poster, filename: Optional[str] = None) -> str:
        filename = filename or os.path.join(self.output_dir, "poster_preview.md")
        md = self.poster_to_markdown(poster)
        with open(filename, "w", encoding="utf-8") as fh:
            fh.write(md)
        logger.info("Markdown preview written to %s", filename)
        return filename

    def poster_to_pptx(self, poster: Poster, out_path: Optional[str] = None) -> str:
        """
        Create a simple PPTX with one slide per section (title + bullets + images).
        Requires python-pptx. This is a simple layout; for production-level poster layouts,
        integrate with a proper layout engine or use templating.
        """
        if not PPTX_AVAILABLE:
            raise RuntimeError("python-pptx is not installed. Install with: pip install python-pptx")

        out_path = out_path or os.path.join(self.output_dir, "poster.pptx")
        prs = Presentation()
        # use a blank slide layout (index 6 may vary)
        blank_layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[5]

        for sec in poster.sections:
            slide = prs.slides.add_slide(blank_layout)
            # Title box
            left = Inches(0.5)
            top = Inches(0.2)
            width = Inches(9)
            height = Inches(1)
            title_box = slide.shapes.add_textbox(left, top, width, height)
            tf = title_box.text_frame
            p = tf.paragraphs[0]
            p.text = sec.title
            p.font.size = Pt(28)
            # Bullets box
            left = Inches(0.5)
            top = Inches(1.2)
            width = Inches(5.5)
            height = Inches(4)
            tb = slide.shapes.add_textbox(left, top, width, height)
            t = tb.text_frame
            bullets = sec.metadata.get("bullets") or []
            for i, b in enumerate(bullets):
                if i == 0:
                    t.text = b
                else:
                    p = t.add_paragraph()
                    p.text = b
                    p.level = 1

            # Add images to the right column (if paths exist)
            img_left = Inches(6.2)
            img_top = Inches(1.2)
            img_w = Inches(3.2)
            img_h = Inches(3.0)
            for i, img in enumerate(sec.images):
                img_path = img.path
                if img_path and os.path.isfile(img_path):
                    try:
                        slide.shapes.add_picture(img_path, img_left, img_top + Inches(i * 3.2), width=img_w, height=img_h)
                    except Exception as e:
                        logger.debug("Failed to add image to pptx: %s", e)

        prs.save(out_path)
        logger.info("PPTX poster written to %s", out_path)
        return out_path
